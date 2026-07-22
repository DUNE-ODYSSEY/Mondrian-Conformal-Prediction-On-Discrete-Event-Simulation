"""
Full project results deck (through Week 15, plus the publication-rigor
upgrade: 30-repeat statistical significance, CQR/Mondrian-CQR, a second
surrogate architecture, a second department) - not the final end-sem PPT,
which will follow the instructor's specific guidelines once shared. This is
a comprehensive working deck: all results in order, why Mondrian CP was
used and why it's better, what was accomplished, and an explicit verdict on
whether the project's novelty target was met.

Re-run any time results/tables/ changes:
.venv\\Scripts\\python.exe slides\\build_full_results_ppt.py

Design: same polished/custom-drawn style as the first mid-sem PPT version -
dark hero title slide, colored kicker/title/underline header on every
content slide, banded tables, card-style callouts.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

OUT_PATH = "slides/full_project_results.pptx"

TEAM = [
    ("G Venugopalan", "CB.AI.U4AID25115"),
    ("Vipin Sudhakar", "CB.AI.U4AID25166"),
    ("Rithvik Arulprakash", "CB.AI.U4AID25148"),
    ("Harshith Kv", "CB.AI.U4AID25119"),
]
COURSE_CODE = "23AID201"
PROJECT_TITLE = "Mondrian Conformal Prediction for Uncertainty\nQuantification in ER Discrete-Event Simulation Surrogates"

DARK = RGBColor(0x1B, 0x22, 0x38)
ACCENT = RGBColor(0x3B, 0x82, 0xF6)
ACCENT_DARK = RGBColor(0x1E, 0x40, 0xAF)
TEXT_DARK = RGBColor(0x1F, 0x29, 0x37)
TEXT_MUTED = RGBColor(0x6B, 0x72, 0x80)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_BG = RGBColor(0xF8, 0xFA, 0xFC)
BORDER = RGBColor(0xE2, 0xE8, 0xF0)
GOOD = RGBColor(0x16, 0xA3, 0x4A)
WARN = RGBColor(0xD9, 0x77, 0x06)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
TOTAL_SLIDES = 21


def new_deck():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


def blank_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def add_rect(slide, x, y, w, h, color, line=False):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    if line:
        shape.line.color.rgb = BORDER
        shape.line.width = Pt(0.75)
    else:
        shape.line.fill.background()
    shape.shadow.inherit = False
    return shape


def add_text(slide, x, y, w, h, text, size=18, bold=False, color=TEXT_DARK,
             align=PP_ALIGN.LEFT, font="Calibri", anchor=MSO_ANCHOR.TOP, italic=False):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    lines = text.split("\n")
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = line
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.italic = italic
        run.font.color.rgb = color
        run.font.name = font
    return box


def add_bullets(slide, x, y, w, h, items, size=16, color=TEXT_DARK, space_after=10, font="Calibri"):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(space_after)
        text, level = item if isinstance(item, tuple) else (item, 0)
        p.level = level
        bullet = "•  " if level == 0 else "‒  "
        run = p.add_run()
        run.text = bullet + text
        run.font.size = Pt(size - level * 2)
        run.font.color.rgb = color
        run.font.name = font
    return box


def slide_header(slide, kicker, title, num, total=TOTAL_SLIDES, title_size=26):
    add_rect(slide, 0, 0, SLIDE_W, Inches(0.09), ACCENT)
    add_text(slide, Inches(0.6), Inches(0.3), Inches(11.5), Inches(0.35),
              kicker.upper(), size=13, bold=True, color=ACCENT, font="Calibri")
    add_text(slide, Inches(0.6), Inches(0.62), Inches(11.9), Inches(0.7),
              title, size=title_size, bold=True, color=DARK, font="Calibri")
    add_rect(slide, Inches(0.6), Inches(1.28), Inches(1.1), Pt(3), ACCENT)
    add_text(slide, Inches(12.2), Inches(7.05), Inches(0.9), Inches(0.35),
              f"{num} / {total}", size=11, color=TEXT_MUTED, align=PP_ALIGN.RIGHT)


def add_table(slide, x, y, w, h, headers, rows, col_widths=None, font_size=13,
              header_color=ACCENT_DARK, highlight_map=None):
    n_rows = len(rows) + 1
    n_cols = len(headers)
    gtable = slide.shapes.add_table(n_rows, n_cols, x, y, w, h).table
    if col_widths:
        for i, cw in enumerate(col_widths):
            gtable.columns[i].width = cw
    for j, hdr in enumerate(headers):
        cell = gtable.cell(0, j)
        cell.text = hdr
        cell.fill.solid()
        cell.fill.fore_color.rgb = header_color
        # cell.text splits on "\n" into separate paragraphs, not a single run
        # with a line break — style every paragraph/run, not just the first,
        # or unstyled lines fall back to a much larger default font size.
        for p in cell.text_frame.paragraphs:
            p.alignment = PP_ALIGN.CENTER if j > 0 else PP_ALIGN.LEFT
            for run in p.runs:
                run.font.size = Pt(font_size)
                run.font.bold = True
                run.font.color.rgb = WHITE
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        cell.margin_top = Pt(4)
        cell.margin_bottom = Pt(4)
    for i, row in enumerate(rows):
        for j, val in enumerate(row):
            cell = gtable.cell(i + 1, j)
            cell.text = str(val)
            cell.fill.solid()
            cell.fill.fore_color.rgb = LIGHT_BG if i % 2 == 0 else WHITE
            for p in cell.text_frame.paragraphs:
                p.alignment = PP_ALIGN.CENTER if j > 0 else PP_ALIGN.LEFT
                for run in p.runs:
                    run.font.size = Pt(font_size)
                    run.font.color.rgb = TEXT_DARK
                    if highlight_map and (i, j) in highlight_map:
                        run.font.color.rgb = highlight_map[(i, j)]
                        run.font.bold = True
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            cell.margin_top = Pt(4)
            cell.margin_bottom = Pt(4)
    return gtable


def pipeline_box(slide, x, y, w, h, label, sublabel, color):
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    box.adjustments[0] = 0.08
    box.fill.solid()
    box.fill.fore_color.rgb = color
    box.line.fill.background()
    box.shadow.inherit = False
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Pt(6)
    tf.margin_right = Pt(6)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r1 = p.add_run()
    r1.text = label
    r1.font.size = Pt(15)
    r1.font.bold = True
    r1.font.color.rgb = WHITE
    if sublabel:
        p2 = tf.add_paragraph()
        p2.alignment = PP_ALIGN.CENTER
        r2 = p2.add_run()
        r2.text = sublabel
        r2.font.size = Pt(10.5)
        r2.font.color.rgb = WHITE
    return box


def pipeline_arrow(slide, x, y, w, h):
    arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, x, y, w, h)
    arrow.fill.solid()
    arrow.fill.fore_color.rgb = TEXT_MUTED
    arrow.line.fill.background()
    arrow.shadow.inherit = False
    return arrow


def build():
    prs = new_deck()
    n = [0]

    def next_num():
        n[0] += 1
        return n[0]

    # ---------------- Slide 1: Title ----------------
    s = blank_slide(prs)
    add_rect(s, 0, 0, SLIDE_W, SLIDE_H, DARK)
    add_rect(s, 0, Inches(6.55), SLIDE_W, Inches(0.12), ACCENT)
    add_text(s, Inches(1), Inches(1.3), Inches(11.3), Inches(2.0),
              PROJECT_TITLE, size=32, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, Inches(1), Inches(3.35), Inches(11.3), Inches(0.5),
              "Full Project Results + Publication-Rigor Upgrade", size=17, color=RGBColor(0x93, 0xC5, 0xFD),
              align=PP_ALIGN.CENTER, italic=True)
    team_text = "   |   ".join([f"{nm} ({i})" for nm, i in TEAM])
    add_text(s, Inches(1), Inches(5.35), Inches(11.3), Inches(0.8),
              team_text, size=13, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, Inches(1), Inches(5.85), Inches(11.3), Inches(0.4),
              f"Course: {COURSE_CODE}", size=13, color=RGBColor(0x93, 0xC5, 0xFD), align=PP_ALIGN.CENTER)
    next_num()

    # ---------------- Slide 2: Problem Statement ----------------
    s = blank_slide(prs)
    slide_header(s, "Motivation", "Problem Statement", next_num())
    add_bullets(s, Inches(0.7), Inches(1.75), Inches(11.9), Inches(5), [
        "Surrogate models are fast, learned approximations of expensive simulations, increasingly used "
        "to guide decisions in stochastic service systems - ER staffing, queueing networks, etc.",
        "A point prediction alone isn't enough for a decision like \"will adding 2 more doctors bring wait "
        "time down enough?\" - that needs a calibrated interval, not just a number.",
        "Gaussian Processes are the established way to get one, but are expensive to train and rely on "
        "distributional assumptions rather than a finite-sample guarantee.",
        "Conformal Prediction (CP) is distribution-free with a finite-sample coverage guarantee and is "
        "cheap to calibrate - but it's only been validated for surrogate models in a narrow set of "
        "domains so far.",
    ], size=18, space_after=16)

    # ---------------- Slide 3: Research Gap ----------------
    s = blank_slide(prs)
    slide_header(s, "Literature", "Research Gap", next_num())
    add_text(s, Inches(0.7), Inches(1.65), Inches(11.9), Inches(0.5),
              "Gopakumar et al. (2026) validate Conformal Prediction for surrogate-model UQ —", size=18, bold=True)
    add_text(s, Inches(0.7), Inches(2.1), Inches(11.9), Inches(0.5),
              "but only in physics domains: PDEs, magnetohydrodynamics (MHD), weather, fusion.", size=18, color=ACCENT_DARK, bold=True)
    add_rect(s, Inches(0.7), Inches(2.95), Inches(11.9), Inches(2.0), LIGHT_BG, line=True)
    add_text(s, Inches(1.0), Inches(3.15), Inches(11.3), Inches(0.4),
              "The paper states two explicit limitations:", size=15, bold=True, color=TEXT_DARK)
    add_bullets(s, Inches(1.0), Inches(3.65), Inches(11.2), Inches(1.2), [
        "Marginal coverage only — no guarantee of coverage within specific subgroups/conditions.",
        "Relies on an exchangeability assumption between calibration and test data — untested under "
        "distribution shift.",
    ], size=15, space_after=10)
    add_text(s, Inches(0.7), Inches(5.25), Inches(11.9), Inches(1.4),
              "Neither limitation has been tested outside physics simulation. Discrete-event / queueing "
              "systems are a fundamentally different domain — stochastic discrete arrivals and departures, "
              "resource contention, priority scheduling — not continuous PDE fields. This is the gap this "
              "project sits in.", size=16, color=TEXT_DARK)

    # ---------------- Slide 4: Why Mondrian CP ----------------
    s = blank_slide(prs)
    slide_header(s, "Approach", "Why Mondrian CP", next_num())
    col_w = Inches(5.7)
    add_rect(s, Inches(0.7), Inches(1.65), col_w, Inches(4.8), WHITE, line=True)
    add_rect(s, Inches(0.7), Inches(1.65), col_w, Inches(0.55), ACCENT_DARK)
    add_text(s, Inches(0.95), Inches(1.77), col_w - Inches(0.5), Inches(0.4),
              "The problem with standard CP", size=15, bold=True, color=WHITE)
    add_bullets(s, Inches(0.95), Inches(2.4), col_w - Inches(0.5), Inches(3.9), [
        "Standard CP calibrates one pooled quantile across every scenario in the calibration set — only "
        "guarantees coverage on average.",
        "A pooled interval is calibrated on average difficulty, so it can badly fail exactly the "
        "scenarios that matter most (understaffed + high demand), while wasting width on easy ones "
        "(overstaffed + low demand).",
        "Averaged coverage can look fine while specific operating conditions quietly fail.",
    ], size=14, space_after=14)

    add_rect(s, Inches(6.7), Inches(1.65), col_w, Inches(4.8), WHITE, line=True)
    add_rect(s, Inches(6.7), Inches(1.65), col_w, Inches(0.55), ACCENT_DARK)
    add_text(s, Inches(6.95), Inches(1.77), col_w - Inches(0.5), Inches(0.4),
              "What Mondrian CP does about it", size=15, bold=True, color=WHITE)
    add_bullets(s, Inches(6.95), Inches(2.4), col_w - Inches(0.5), Inches(3.9), [
        "Calibrates a separate quantile per category — here, staffing tercile x arrival-rate tercile — "
        "instead of one pooled quantile.",
        "Each category gets its own, individually honest coverage guarantee instead of borrowing "
        "strength from categories that don't resemble it.",
        "Directly targets Gopakumar et al.'s marginal-coverage limitation — the first place it's been "
        "tested outside physics simulation.",
    ], size=14, space_after=14)

    # ---------------- Slide 5: Methodology ----------------
    s = blank_slide(prs)
    slide_header(s, "Pipeline", "Methodology Overview", next_num())
    y = Inches(2.7)
    h = Inches(1.5)
    bw = Inches(3.15)
    aw = Inches(0.55)
    x = Inches(0.55)
    pipeline_box(s, x, y, bw, h, "Discrete-Event\nSimulation", "SimPy, calibrated on real ED data", ACCENT_DARK)
    x2 = x + bw
    pipeline_arrow(s, x2, y + h / 2 - Inches(0.2), aw, Inches(0.4))
    x3 = x2 + aw
    pipeline_box(s, x3, y, bw, h, "Surrogate Model", "Gradient boosting, trained on\nDES scenario sweep", ACCENT)
    x4 = x3 + bw
    pipeline_arrow(s, x4, y + h / 2 - Inches(0.2), aw, Inches(0.4))
    x5 = x4 + aw
    pipeline_box(s, x5, y, bw, h, "Uncertainty\nQuantification", "GP baseline vs. Standard CP\nvs. Mondrian CP", RGBColor(0x0F, 0x76, 0x6E))

    add_text(s, Inches(0.55), Inches(4.55), Inches(12.2), Inches(0.4),
              "Each stage validated against the previous before moving on:", size=14, italic=True, color=TEXT_MUTED)
    add_bullets(s, Inches(0.7), Inches(5.0), Inches(11.7), Inches(2.2), [
        "DES output validated against real aggregated ED statistics (91.0% match to real daily volume).",
        "Surrogate accuracy validated against DES outputs held out from training (MAE / RMSE / R²).",
        "GP baseline, standard CP, Mondrian CP all evaluated on the same fixed test split and target "
        "coverage (α = 0.1), so every method is directly comparable.",
        "Standard/Mondrian CP calibrated on a separate 1200-scenario set — never the surrogate's own "
        "training data, since reusing training residuals would understate the true residual spread.",
    ], size=14, space_after=8)

    # ---------------- Slide 6: Real-world data ----------------
    s = blank_slide(prs)
    slide_header(s, "Data", "Real-World Data Used", next_num())
    add_text(s, Inches(0.7), Inches(1.6), Inches(11.9), Inches(0.5),
              "Hospital Triage and Patient History Data (Kaggle, maalona)", size=18, bold=True, color=DARK)
    add_text(s, Inches(0.7), Inches(2.05), Inches(11.9), Inches(0.5),
              "Yale New Haven Health System, retrospective study, March 2014 – July 2017",
              size=14, color=TEXT_MUTED, italic=True)
    stats = [("560,486", "total ED visits"), ("972", "variables per visit"),
             ("3", "EDs (1 academic + 2 community)"), ("258.2 / day", "calibrated rate, Dept. A")]
    sx = Inches(0.7)
    sw = Inches(2.85)
    for i, (num, lbl) in enumerate(stats):
        cx = sx + i * (sw + Inches(0.1))
        add_rect(s, cx, Inches(2.75), sw, Inches(1.15), LIGHT_BG, line=True)
        add_text(s, cx, Inches(2.85), sw, Inches(0.55), num, size=22, bold=True, color=ACCENT_DARK, align=PP_ALIGN.CENTER)
        add_text(s, cx, Inches(3.4), sw, Inches(0.4), lbl, size=11, color=TEXT_MUTED, align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.7), Inches(4.25), Inches(11.9), Inches(0.4),
              "What we extracted — real vs. literature-calibrated (disclosed explicitly, not blended silently):",
              size=15, bold=True, color=TEXT_DARK)
    add_rect(s, Inches(0.7), Inches(4.75), Inches(5.8), Inches(2.0), WHITE, line=True)
    add_text(s, Inches(0.9), Inches(4.85), Inches(5.4), Inches(0.35), "From real data", size=13, bold=True, color=GOOD)
    add_bullets(s, Inches(0.9), Inches(5.2), Inches(5.4), Inches(1.5), [
        "Arrival pattern by 4-hour bin, day of week, month",
        "ESI acuity mix (triage severity distribution)",
        "Department-level daily visit rate",
    ], size=13, space_after=6)
    add_rect(s, Inches(6.65), Inches(4.75), Inches(5.95), Inches(2.0), WHITE, line=True)
    add_text(s, Inches(6.85), Inches(4.85), Inches(5.4), Inches(0.35), "From literature (not in dataset)", size=13, bold=True, color=WARN)
    add_bullets(s, Inches(6.85), Inches(5.2), Inches(5.4), Inches(1.5), [
        "Service/treatment time per ESI level (log-normal)",
        "Dataset has no discharge timestamp — confirmed by",
        ("exhaustive search across all 972 columns", 1),
    ], size=13, space_after=6)

    # ---------------- Slide 7: DES Validation ----------------
    s = blank_slide(prs)
    slide_header(s, "Results", "DES Validation", next_num())
    add_table(s, Inches(0.7), Inches(1.8), Inches(11.9), Inches(1.7),
              ["Metric (200 simulated days)", "Value"],
              [
                  ["Real visits/day (Dept. A)", "258.2"],
                  ["Simulated mean visits/day", "235.1"],
                  ["Match", "91.0%"],
              ], col_widths=[Inches(8.5), Inches(3.4)], font_size=15)
    add_rect(s, Inches(0.7), Inches(3.9), Inches(11.9), Inches(1.7), LIGHT_BG, line=True)
    add_text(s, Inches(1.0), Inches(4.1), Inches(11.3), Inches(1.3),
              "The ~9% shortfall is expected, not a calibration flaw: patients still queued when a "
              "simulated 24h day ends are right-censored out of that day's stats rather than carried into "
              "a \"next day,\" since each run is meant to be one independent sample for surrogate training, "
              "not a continuous rollout. This censoring effect resurfaces later in the exchangeability "
              "stress test (Slide 13).", size=14, color=TEXT_DARK)

    # ---------------- Slide 8: Surrogate accuracy ----------------
    s = blank_slide(prs)
    slide_header(s, "Results", "Surrogate Model Accuracy", next_num())
    add_table(s, Inches(0.7), Inches(1.8), Inches(11.9), Inches(2.2),
              ["Target", "MAE", "RMSE", "R²"],
              [
                  ["n_patients", "9.90", "12.59", "0.929"],
                  ["mean_wait_minutes", "8.86", "13.23", "0.787"],
                  ["mean_total_minutes", "9.88", "13.61", "0.762"],
                  ["p95_wait_minutes", "66.94", "102.47", "0.647"],
              ], col_widths=[Inches(4.5), Inches(2.5), Inches(2.5), Inches(2.4)], font_size=14)
    add_rect(s, Inches(0.7), Inches(4.3), Inches(11.9), Inches(1.9), LIGHT_BG, line=True)
    add_text(s, Inches(1.0), Inches(4.5), Inches(11.3), Inches(1.5),
              "Gradient boosting (HistGradientBoostingRegressor), trained on 5000 DES scenarios sweeping "
              "staffing and arrival-rate combinations. p95_wait_minutes fits worst — a tail statistic from "
              "one noisy simulated day depends heavily on the specific stochastic realization, not just "
              "the two scenario parameters. This matters later in the exchangeability stress test.",
              size=14, color=TEXT_DARK)

    # ---------------- Slide 9: GP baseline ----------------
    s = blank_slide(prs)
    slide_header(s, "Results", "GP Baseline (the benchmark)", next_num())
    add_table(s, Inches(0.7), Inches(1.8), Inches(11.9), Inches(2.2),
              ["Target", "Target coverage", "Empirical coverage", "Mean width"],
              [
                  ["n_patients", "90%", "88.5%", "39.3"],
                  ["mean_wait_minutes", "90%", "87.7%", "43.6"],
                  ["mean_total_minutes", "90%", "88.8%", "43.4"],
                  ["p95_wait_minutes", "90%", "90.1%", "350.3"],
              ], col_widths=[Inches(3.9), Inches(2.7), Inches(2.7), Inches(2.6)], font_size=14)
    add_rect(s, Inches(0.7), Inches(4.3), Inches(11.9), Inches(1.3), LIGHT_BG, line=True)
    add_text(s, Inches(1.0), Inches(4.5), Inches(11.3), Inches(0.9),
              "GP undercovers on 3 of 4 targets — expected, since it relies on distributional assumptions "
              "rather than a finite-sample guarantee. This is the gap the rest of the project tests "
              "whether CP / Mondrian CP can close.", size=14, color=TEXT_DARK)

    # ---------------- Slide 10: Standard CP ----------------
    s = blank_slide(prs)
    slide_header(s, "Results", "Standard Conformal Prediction", next_num())
    add_table(s, Inches(0.7), Inches(1.8), Inches(11.9), Inches(2.2),
              ["Target", "Sym. coverage", "Sym. width", "Asym. coverage", "Asym. width"],
              [
                  ["n_patients", "92.1%", "43.6", "92.0%", "43.6"],
                  ["mean_wait_minutes", "89.0%", "47.2", "88.7%", "47.2"],
                  ["mean_total_minutes", "90.2%", "46.4", "89.6%", "47.2"],
                  ["p95_wait_minutes", "90.8%", "362.9", "89.7%", "355.6"],
              ], col_widths=[Inches(3.3), Inches(2.15), Inches(2.15), Inches(2.15), Inches(2.15)], font_size=13)
    add_rect(s, Inches(0.7), Inches(4.3), Inches(11.9), Inches(1.9), LIGHT_BG, line=True)
    add_text(s, Inches(1.0), Inches(4.5), Inches(11.3), Inches(1.5),
              "Two nonconformity measures tested: symmetric |y − yhat|, and an asymmetric version using "
              "separate upper/lower residual quantiles for skewed targets. CP already lands closer to the "
              "90% target than GP, with no systematic undercoverage bias — but it's still a single pooled "
              "quantile, so it only guarantees coverage on average, not within any specific scenario.",
              size=14, color=TEXT_DARK)

    # ---------------- Slide 11: Mondrian CP core finding ----------------
    s = blank_slide(prs)
    slide_header(s, "Results — Core Finding", "Mondrian CP Closes the Gap", next_num(), title_size=24)
    add_table(s, Inches(0.7), Inches(1.7), Inches(11.9), Inches(1.7),
              ["Target", "Pooled coverage\n(staff=Low / arrival=High)", "Mondrian coverage\n(same category)", "Target"],
              [
                  ["mean_wait_minutes", "68.2%", "90.9%", "90%"],
                  ["mean_total_minutes", "80.7%", "92.0%", "90%"],
                  ["p95_wait_minutes", "72.7%", "92.0%", "90%"],
              ], col_widths=[Inches(3.2), Inches(3.4), Inches(2.9), Inches(2.4)], font_size=13)
    add_rect(s, Inches(0.7), Inches(3.65), Inches(11.9), Inches(3.55), LIGHT_BG, line=True)
    add_text(s, Inches(1.0), Inches(3.85), Inches(11.3), Inches(3.2),
              "The worst pooled-CP category is consistently staff=Low/arrival=High — understaffed during "
              "a demand surge, the scenario that matters most operationally. A pooled quantile is "
              "calibrated on average difficulty, so it silently fails exactly this scenario while "
              "overcovering the easy one (staff=High/arrival=Low hits 100% pooled coverage — wasted "
              "width). Mondrian gives the hard category an honestly wider interval instead of pretending "
              "it's as easy as the rest.\n\n"
              "n_patients is the one target where Mondrian does not help — its pooled coverage was already "
              "fairly uniform across categories (87–99%), so there was no real conditional gap to fix, and "
              "splitting into 9 smaller calibration cells just adds finite-sample noise. That's an honest, "
              "expected Mondrian CP tradeoff, not a failure.",
              size=14, color=TEXT_DARK)

    # ---------------- Slide 12: Full comparison ----------------
    s = blank_slide(prs)
    slide_header(s, "Results", "GP vs. Standard CP vs. Mondrian CP", next_num())
    s.shapes.add_picture("results/figures/full_comparison.png", Inches(0.7), Inches(1.55), width=Inches(11.9))
    add_text(s, Inches(0.7), Inches(6.45), Inches(11.9), Inches(0.9),
              "CP calibrates ~650–1000x faster than GP (~0.01s vs. 7–11.5s). Mondrian CP's marginal "
              "coverage matches or beats standard CP's despite ~9x smaller per-category calibration "
              "slices, and is also usually narrower — pooling was wasting width in easy categories. "
              "(Single-split result — see the 30-repeat validated version on the next slide.)",
              size=13, italic=True, color=TEXT_MUTED)

    # ---------------- Slide 12b: Publication-grade comparison (5 methods, 30 repeats) ----------------
    s = blank_slide(prs)
    slide_header(s, "Results — Validated", "All 5 Methods, 30 Independent Repeats", next_num(), title_size=24)
    s.shapes.add_picture("results/figures/publication_comparison.png", Inches(0.7), Inches(1.5), width=Inches(11.9))
    add_text(s, Inches(0.7), Inches(6.6), Inches(11.9), Inches(0.8),
              "Error bars are 95% CI on the mean across 30 independent calibration/test draws — tight "
              "enough to read as single dots, meaning these differences are stable, not single-split noise. "
              "CQR (adaptive, quantile-regression-based) joins the comparison as a fifth method.",
              size=13, italic=True, color=TEXT_MUTED)

    # ---------------- Slide 13: Exchangeability stress test ----------------
    s = blank_slide(prs)
    slide_header(s, "Results", "Exchangeability Stress Test", next_num())
    add_table(s, Inches(0.7), Inches(1.6), Inches(11.9), Inches(2.0),
              ["Arrival multiplier", "n_patients", "mean_wait", "mean_total", "p95_wait"],
              [
                  ["1.3 (in range)", "93.3%", "83.0%", "88.0%", "82.0%"],
                  ["1.5", "78.7%", "68.7%", "73.7%", "69.7%"],
                  ["1.8", "70.3%", "42.7%", "47.3%", "43.0%"],
                  ["2.0", "64.7%", "34.3%", "39.0%", "33.0%"],
                  ["3.0", "31.7%", "12.3%", "4.7%", "78.3%*"],
              ], col_widths=[Inches(2.3), Inches(2.4), Inches(2.4), Inches(2.4), Inches(2.4)], font_size=13)
    add_rect(s, Inches(0.7), Inches(3.85), Inches(11.9), Inches(3.35), LIGHT_BG, line=True)
    add_text(s, Inches(1.0), Inches(4.0), Inches(11.3), Inches(3.1),
              "Standard CP shown (Mondrian tracks closely). Root cause, verified directly: the gradient-"
              "boosting surrogate cannot extrapolate past its training range — mean_wait_minutes's "
              "prediction is frozen at 42.4 for every multiplier from 1.3 to 3.0 while the true value "
              "keeps climbing.\n\n"
              "*p95_wait's apparent recovery at 3.0x isn't the interval working — it's the Slide 7 "
              "day-boundary censoring effect getting worse at extreme overload (fewer patients finish "
              "within a day, shrinking the completed-visit pool the p95 is computed over).\n\n"
              "Both stated Gopakumar et al. limitations are now tested in this domain.",
              size=13, color=TEXT_DARK)

    # ---------------- Slide 14: Statistical rigor (30 repeats) ----------------
    s = blank_slide(prs)
    slide_header(s, "Publication-Rigor Upgrade", "Statistical Rigor: 30 Independent Repeats", next_num(), title_size=22)
    add_text(s, Inches(0.7), Inches(1.6), Inches(11.9), Inches(0.9),
              "Every result so far was a single train/test split — no way to tell if a number like "
              "\"68.2% coverage\" was a robust effect or one unlucky draw. Repeated the full GP / "
              "standard-CP / Mondrian-CP pipeline across 30 independent (calibration, test) draws.",
              size=15, color=TEXT_DARK)
    add_table(s, Inches(0.7), Inches(2.65), Inches(11.9), Inches(2.4),
              ["Target", "Method", "Coverage (mean ± 95% CI)", "Width (mean ± 95% CI)"],
              [
                  ["mean_wait_minutes", "GP", "88.31% ± 0.47%", "44.4 ± 1.1"],
                  ["mean_wait_minutes", "Standard CP", "90.09% ± 0.46%", "48.0 ± 1.5"],
                  ["mean_wait_minutes", "Mondrian CP", "91.07% ± 0.35%", "41.7 ± 1.4"],
                  ["p95_wait_minutes", "GP", "88.97% ± 0.44%", "354.6 ± 11.5"],
                  ["p95_wait_minutes", "Standard CP", "90.06% ± 0.39%", "377.1 ± 12.6"],
                  ["p95_wait_minutes", "Mondrian CP", "91.35% ± 0.46%", "328.6 ± 13.4"],
              ], col_widths=[Inches(2.9), Inches(2.4), Inches(3.4), Inches(3.2)], font_size=12)
    add_rect(s, Inches(0.7), Inches(5.3), Inches(11.9), Inches(1.75), LIGHT_BG, line=True)
    add_text(s, Inches(1.0), Inches(5.45), Inches(11.3), Inches(1.5),
              "Paired t-test (same 30 draws underlie all methods): Mondrian CP's coverage advantage over "
              "both standard CP and GP is significant at p < 0.001 on every one of the 4 targets, most "
              "p-values below 1e-05. With averaging, Mondrian CP's coverage now sits consistently above "
              "90% while GP consistently undercovers — a statistically defensible claim, not just a "
              "directionally suggestive one.", size=13, color=TEXT_DARK)

    # ---------------- Slide 15: CQR ----------------
    s = blank_slide(prs)
    slide_header(s, "Publication-Rigor Upgrade", "CQR: A Stronger, Adaptive Baseline", next_num(), title_size=24)
    add_text(s, Inches(0.7), Inches(1.6), Inches(11.9), Inches(1.0),
              "Conformalized Quantile Regression (CQR) calibrates an interval that adapts to the local "
              "spread predicted by quantile regressors, instead of one fixed width (standard CP) or one "
              "fixed width per category (Mondrian CP) built from a mean predictor.",
              size=15, color=TEXT_DARK)
    add_table(s, Inches(0.7), Inches(2.7), Inches(11.9), Inches(1.7),
              ["Target", "Standard CP", "Mondrian CP", "CQR", "Mondrian CQR"],
              [
                  ["p95_wait_minutes (coverage)", "90.1%", "91.4%", "90.9%", "92.2%"],
                  ["p95_wait_minutes (width)", "377.1", "328.6", "275.6", "292.7"],
              ], col_widths=[Inches(3.5), Inches(2.1), Inches(2.1), Inches(2.1), Inches(2.1)], font_size=12)
    add_rect(s, Inches(0.7), Inches(4.7), Inches(11.9), Inches(2.35), LIGHT_BG, line=True)
    add_text(s, Inches(1.0), Inches(4.85), Inches(11.3), Inches(2.1),
              "CQR dominates standard CP outright on p95_wait_minutes — equal-or-better coverage, ~27% "
              "narrower interval (p < 0.005). Against Mondrian CP it's a genuine tradeoff: slightly lower "
              "coverage, but narrower — different points on the same coverage/width frontier. Combining "
              "both ideas (Mondrian CQR) wins on both axes simultaneously for p95_wait_minutes — higher "
              "coverage AND narrower than Mondrian CP alone. The five methods form a real frontier, not a "
              "strict ranking: each adds value in a different way.", size=13, color=TEXT_DARK)

    # ---------------- Slide 16: MLP robustness check ----------------
    s = blank_slide(prs)
    slide_header(s, "Publication-Rigor Upgrade", "Robustness Check: A Second Surrogate", next_num(), title_size=24)
    add_text(s, Inches(0.7), Inches(1.6), Inches(11.9), Inches(0.9),
              "Trained an MLP (neural network) surrogate — nearly identical in-distribution accuracy to "
              "gradient boosting (R² within ~0.01 on every target) — and re-ran the exchangeability stress "
              "test to see if the coverage collapse was specific to tree-based extrapolation.",
              size=15, color=TEXT_DARK)
    add_table(s, Inches(0.7), Inches(2.6), Inches(11.9), Inches(1.4),
              ["n_capacity=30", "Arrival 1.0x", "Arrival 1.3x", "Arrival 2.0x", "Arrival 3.0x", "True (3.0x)"],
              [
                  ["GBR prediction", "240.5", "247.8", "247.8 (frozen)", "247.8 (frozen)", "280.9"],
                  ["MLP prediction", "234.8", "245.8", "336.5", "521.2", "280.9"],
              ], col_widths=[Inches(2.3), Inches(1.9), Inches(1.9), Inches(2.0), Inches(1.9), Inches(1.9)], font_size=11)
    add_rect(s, Inches(0.7), Inches(4.25), Inches(11.9), Inches(2.8), LIGHT_BG, line=True)
    add_text(s, Inches(1.0), Inches(4.4), Inches(11.3), Inches(2.55),
              "Counter-intuitive result: the true DES output saturates at extreme overload (the same "
              "day-boundary censoring effect from Slide 7 — fewer patients finish within a day). GBR's "
              "frozen prediction happens to land close to that saturating value by coincidence (error "
              "~33). The MLP keeps extrapolating the upward trend it learned near the training boundary "
              "and overshoots massively (error ~240, 7x worse) — n_patients and mean_total_minutes both "
              "hit exactly 0% coverage with the MLP by 2.0x severity, worse than GBR's 31.7%/4.7% at "
              "3.0x. The ability to extrapolate is not automatically better than not being able to — "
              "confidently-wrong extrapolation is worse than frozen-but-plausible extrapolation here. "
              "Strengthens the exchangeability finding: the breakdown isn't one architecture's quirk.",
              size=13, color=TEXT_DARK)

    # ---------------- Slide 17: Second department ----------------
    s = blank_slide(prs)
    slide_header(s, "Publication-Rigor Upgrade", "Generalizability: A Second Department", next_num(), title_size=24)
    add_text(s, Inches(0.7), Inches(1.6), Inches(11.9), Inches(0.9),
              "Recalibrated the entire pipeline — DES, surrogate, CP — on Department B (166,497 visits, "
              "roughly half A's volume, meaningfully different ESI acuity mix — a community rather than "
              "academic site) to test whether the core finding is specific to one department.",
              size=15, color=TEXT_DARK)
    add_table(s, Inches(0.7), Inches(2.65), Inches(11.9), Inches(1.7),
              ["Target", "Dept. B: pooled (worst category)", "Dept. B: Mondrian", "Dept. A (for comparison)"],
              [
                  ["mean_wait_minutes", "76.2%", "89.3%", "68.2% -> 90.9%"],
                  ["mean_total_minutes", "81.0%", "90.5%", "80.7% -> 92.0%"],
                  ["p95_wait_minutes", "81.0%", "86.9%", "72.7% -> 92.0%"],
              ], col_widths=[Inches(3.0), Inches(3.2), Inches(2.5), Inches(3.2)], font_size=12)
    add_rect(s, Inches(0.7), Inches(4.65), Inches(11.9), Inches(2.4), LIGHT_BG, line=True)
    add_text(s, Inches(1.0), Inches(4.8), Inches(11.3), Inches(2.15),
              "Same worst category (understaffed + high demand) in both independent sites, similar "
              "correction magnitude from Mondrian CP, same wasted-width pattern on the easy category. "
              "Honest difference worth keeping: n_patients behaves differently between sites — a real "
              "conditional gap appears in Department B that wasn't present in A, and Mondrian helps there "
              "this time. That's disclosed, not smoothed over — and it's part of what makes the "
              "replication credible rather than suspiciously perfect.", size=13, color=TEXT_DARK)

    # ---------------- Slide 18: What we accomplished ----------------
    s = blank_slide(prs)
    slide_header(s, "Summary", "What We Accomplished", next_num())
    add_bullets(s, Inches(0.7), Inches(1.75), Inches(11.9), Inches(5), [
        "Built and validated an ER discrete-event simulation (SimPy) calibrated on real hospital arrival "
        "data — 91.0% match to real daily patient volume.",
        "Generated a 5000-scenario training sweep and trained a surrogate model with strong accuracy on "
        "3 of 4 targets (R² 0.76–0.93).",
        "Implemented and compared 5 UQ methods — GP baseline, standard CP, Mondrian CP, CQR, and "
        "Mondrian CQR — all on the same test data and target coverage.",
        "Demonstrated Mondrian CP closing a real marginal-vs-conditional coverage gap (68–81% pooled "
        "coverage in the hardest category, corrected to 91–92%), validated with paired significance "
        "tests (p < 0.001) across 30 independent repeats, not a single lucky split.",
        "Ran an exchangeability stress test across 8 severity levels with two structurally different "
        "surrogate architectures, and root-caused the coverage collapse to a specific, verified "
        "mechanism in each case, not just \"it broke.\"",
        "Replicated the core finding at a second, independent department with different volume and "
        "acuity mix — the marginal-vs-conditional gap Mondrian CP closes is not a one-site artifact.",
    ], size=15, space_after=11)

    # ---------------- Slide 19: Did we meet our novelty ----------------
    s = blank_slide(prs)
    slide_header(s, "Verdict", "Did We Meet Our Novelty Target?", next_num(), title_size=24)
    add_text(s, Inches(0.7), Inches(1.7), Inches(11.9), Inches(1.1),
              "The claim was to test Gopakumar et al. (2026)'s two stated CP limitations — marginal "
              "coverage, exchangeability — outside the physics domains they validated CP in, using a "
              "discrete-event / queueing simulation surrogate instead.", size=16, color=TEXT_DARK)
    add_rect(s, Inches(0.7), Inches(3.0), Inches(11.9), Inches(2.5), LIGHT_BG, line=True)
    add_bullets(s, Inches(1.0), Inches(3.2), Inches(11.3), Inches(2.2), [
        "Marginal coverage limitation: tested, and now statistically validated (p < 0.001, 30 repeats) "
        "and replicated at a second, independent site.",
        "Exchangeability limitation: tested with two structurally different surrogate architectures — "
        "confirmed to break down outside the training distribution regardless of architecture, with a "
        "specific, verified root cause in each case rather than a generic \"performance degrades\" claim.",
    ], size=15, space_after=10)
    add_text(s, Inches(0.7), Inches(5.75), Inches(11.9), Inches(1.2),
              "Yes — both limitations were tested in a domain they hadn't been tested in before, with "
              "concrete, statistically validated, cross-architecture and cross-site answers, not just a "
              "single replication of the base paper's physics-domain results.", size=16, bold=True, color=GOOD)

    # ---------------- Slide 20: Limitations & what's left ----------------
    s = blank_slide(prs)
    slide_header(s, "Honesty Check", "Limitations & What's Left", next_num())
    add_bullets(s, Inches(0.7), Inches(1.75), Inches(11.9), Inches(5), [
        "Service/treatment time is literature-calibrated, not derived from the dataset — the dataset "
        "has no discharge timestamp. Disclosed explicitly throughout, not blended silently into the "
        "\"real data\" story.",
        "Generalizability tested across 2 of the dataset's 3 departments (one academic, one community) — "
        "both replicate the core finding, but the third department and other health systems entirely are "
        "still untested.",
        "The exchangeability breakdown was tested with 2 surrogate architectures (gradient boosting, "
        "MLP) — both fail outside the training range, via different mechanisms. Other architecture "
        "families (e.g. a GP-as-surrogate) are untested.",
        "Literature review (30 papers, 3 in depth) still in progress.",
        "This is a working results deck, not the final end-sem PPT — that will follow the instructor's "
        "specific guidelines once shared.",
    ], size=16, space_after=14)

    prs.save(OUT_PATH)
    print(f"Saved {OUT_PATH}")


if __name__ == "__main__":
    build()
