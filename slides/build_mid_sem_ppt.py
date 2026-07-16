"""
Builds the mid-sem review PPT from the real results in results/tables/.
Re-run any time those tables change: .venv\\Scripts\\python.exe slides\\build_mid_sem_ppt.py

Slide 8 (literature review snapshot) is left as a placeholder — the
literature review is being done separately and isn't finished yet, so
faking category/citation counts here would misrepresent progress.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

OUT_PATH = "slides/mid_sem_presentation.pptx"

TEAM = [
    ("G Venugopalan", "CB.AI.U4AID25115"),
    ("Vipin Sudhakar", "CB.AI.U4AID25166"),
    ("Rithvik Arulprakash", "CB.AI.U4AID25148"),
    ("Harshith Kv", "CB.AI.U4AID25119"),
]
COURSE_CODE = "23AID201"
PROJECT_TITLE = "Mondrian Conformal Prediction for Uncertainty\nQuantification in ER Discrete-Event Simulation Surrogates"

# Palette
DARK = RGBColor(0x1B, 0x22, 0x38)
ACCENT = RGBColor(0x3B, 0x82, 0xF6)
ACCENT_DARK = RGBColor(0x1E, 0x40, 0xAF)
TEXT_DARK = RGBColor(0x1F, 0x29, 0x37)
TEXT_MUTED = RGBColor(0x6B, 0x72, 0x80)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_BG = RGBColor(0xF8, 0xFA, 0xFC)
BORDER = RGBColor(0xE2, 0xE8, 0xF0)
GOOD = RGBColor(0x16, 0xA3, 0x4A)
WARN = RGBColor(0xD9, 0x77, 0x06)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def new_deck():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


def blank_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def add_rect(slide, x, y, w, h, color, line=False):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    if line:
        shape.line.color.rgb = BORDER
        shape.line.width = Pt(0.75)
    else:
        shape.line.fill.background()
    shape.shadow.inherit = False
    return shape


def add_text(slide, x, y, w, h, text, size=18, bold=False, color=TEXT_DARK,
             align=PP_ALIGN.LEFT, font="Calibri", anchor=MSO_ANCHOR.TOP, italic=False):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    lines = text.split("\n")
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = line
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.italic = italic
        run.font.color.rgb = color
        run.font.name = font
    return box


def add_bullets(slide, x, y, w, h, items, size=16, color=TEXT_DARK, space_after=10, font="Calibri"):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(space_after)
        if isinstance(item, tuple):
            text, level = item
        else:
            text, level = item, 0
        p.level = level
        bullet = "•  " if level == 0 else "‒  "
        run = p.add_run()
        run.text = bullet + text
        run.font.size = Pt(size - level * 2)
        run.font.color.rgb = color
        run.font.name = font
    return box


def slide_header(slide, kicker, title, num, total=9):
    add_rect(slide, 0, 0, SLIDE_W, Inches(0.09), ACCENT)
    add_text(slide, Inches(0.6), Inches(0.35), Inches(8), Inches(0.35),
              kicker.upper(), size=13, bold=True, color=ACCENT, font="Calibri")
    add_text(slide, Inches(0.6), Inches(0.68), Inches(11.5), Inches(0.7),
              title, size=28, bold=True, color=DARK, font="Calibri")
    add_rect(slide, Inches(0.6), Inches(1.35), Inches(1.1), Pt(3), ACCENT)
    add_text(slide, Inches(12.2), Inches(7.05), Inches(0.9), Inches(0.35),
              f"{num} / {total}", size=11, color=TEXT_MUTED, align=PP_ALIGN.RIGHT)


def add_table(slide, x, y, w, h, headers, rows, col_widths=None, font_size=13,
              header_color=ACCENT_DARK, highlight_col=None, highlight_map=None):
    n_rows = len(rows) + 1
    n_cols = len(headers)
    gtable = slide.shapes.add_table(n_rows, n_cols, x, y, w, h).table
    if col_widths:
        for i, cw in enumerate(col_widths):
            gtable.columns[i].width = cw
    for j, hdr in enumerate(headers):
        cell = gtable.cell(0, j)
        cell.text = hdr
        cell.fill.solid()
        cell.fill.fore_color.rgb = header_color
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER if j > 0 else PP_ALIGN.LEFT
        run = p.runs[0]
        run.font.size = Pt(font_size)
        run.font.bold = True
        run.font.color.rgb = WHITE
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        cell.margin_top = Pt(4)
        cell.margin_bottom = Pt(4)
    for i, row in enumerate(rows):
        for j, val in enumerate(row):
            cell = gtable.cell(i + 1, j)
            cell.text = str(val)
            cell.fill.solid()
            cell.fill.fore_color.rgb = LIGHT_BG if i % 2 == 0 else WHITE
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER if j > 0 else PP_ALIGN.LEFT
            run = p.runs[0]
            run.font.size = Pt(font_size)
            run.font.color.rgb = TEXT_DARK
            if highlight_map and (i, j) in highlight_map:
                run.font.color.rgb = highlight_map[(i, j)]
                run.font.bold = True
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            cell.margin_top = Pt(4)
            cell.margin_bottom = Pt(4)
    return gtable


def pipeline_box(slide, x, y, w, h, label, sublabel, color):
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    box.adjustments[0] = 0.08
    box.fill.solid()
    box.fill.fore_color.rgb = color
    box.line.fill.background()
    box.shadow.inherit = False
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Pt(6)
    tf.margin_right = Pt(6)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r1 = p.add_run()
    r1.text = label
    r1.font.size = Pt(15)
    r1.font.bold = True
    r1.font.color.rgb = WHITE
    if sublabel:
        p2 = tf.add_paragraph()
        p2.alignment = PP_ALIGN.CENTER
        r2 = p2.add_run()
        r2.text = sublabel
        r2.font.size = Pt(10.5)
        r2.font.color.rgb = WHITE
    return box


def pipeline_arrow(slide, x, y, w, h):
    arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, x, y, w, h)
    arrow.fill.solid()
    arrow.fill.fore_color.rgb = TEXT_MUTED
    arrow.line.fill.background()
    arrow.shadow.inherit = False
    return arrow


def build():
    prs = new_deck()

    # ---------------- Slide 1: Title ----------------
    s = blank_slide(prs)
    add_rect(s, 0, 0, SLIDE_W, SLIDE_H, DARK)
    add_rect(s, 0, Inches(6.55), SLIDE_W, Inches(0.12), ACCENT)
    add_text(s, Inches(1), Inches(1.5), Inches(11.3), Inches(2.0),
              PROJECT_TITLE, size=34, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, Inches(1), Inches(3.55), Inches(11.3), Inches(0.5),
              "Mid-Semester Review", size=18, color=RGBColor(0x93, 0xC5, 0xFD), align=PP_ALIGN.CENTER, italic=True)

    team_text = "   |   ".join([f"{n} ({i})" for n, i in TEAM])
    add_text(s, Inches(1), Inches(5.35), Inches(11.3), Inches(0.8),
              team_text, size=13, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, Inches(1), Inches(5.85), Inches(11.3), Inches(0.4),
              f"Course: {COURSE_CODE}", size=13, color=RGBColor(0x93, 0xC5, 0xFD), align=PP_ALIGN.CENTER)

    # ---------------- Slide 2: Problem Statement ----------------
    s = blank_slide(prs)
    slide_header(s, "Motivation", "Problem Statement", 2)
    add_bullets(s, Inches(0.7), Inches(1.75), Inches(11.9), Inches(4.5), [
        "Surrogate models — fast learned approximations of expensive simulations — are increasingly used "
        "to guide decisions in stochastic service systems (e.g. ER staffing, queueing networks).",
        "A surrogate's point prediction alone is not enough for high-stakes operational decisions: "
        "“will adding 2 more doctors cut wait time enough?” needs a calibrated interval, not just a number.",
        "Existing UQ approaches have real gaps here:",
        ("Gaussian Processes: well-established, but computationally expensive (cubic training cost) "
         "and rely on distributional assumptions rather than a finite-sample guarantee.", 1),
        ("Conformal Prediction (CP): distribution-free, finite-sample coverage guarantees, and cheap to "
         "calibrate — but its reliability for surrogate models has only been validated in specific domains.", 1),
        "This project asks: does CP's reliability hold when the surrogate is trained on a discrete-event, "
        "queueing-based simulation instead of the domains it's been validated on?",
    ], size=17, space_after=14)

    # ---------------- Slide 3: Research Gap ----------------
    s = blank_slide(prs)
    slide_header(s, "Literature", "Research Gap", 3)
    add_text(s, Inches(0.7), Inches(1.65), Inches(11.9), Inches(0.5),
              "Gopakumar et al. (2026) validate Conformal Prediction for surrogate-model UQ —", size=18, bold=True)
    add_text(s, Inches(0.7), Inches(2.1), Inches(11.9), Inches(0.5),
              "but only in physics domains: PDEs, magnetohydrodynamics (MHD), weather, fusion.", size=18, color=ACCENT_DARK, bold=True)

    add_rect(s, Inches(0.7), Inches(2.95), Inches(11.9), Inches(2.3), LIGHT_BG, line=True)
    add_text(s, Inches(1.0), Inches(3.15), Inches(11.3), Inches(0.4),
              "The paper states two explicit limitations:", size=15, bold=True, color=TEXT_DARK)
    add_bullets(s, Inches(1.0), Inches(3.65), Inches(11.2), Inches(1.4), [
        "Marginal coverage only — no guarantee of coverage within specific subgroups/conditions.",
        "Relies on an exchangeability assumption between calibration and test data — untested under "
        "distribution shift.",
    ], size=15, space_after=10)

    add_text(s, Inches(0.7), Inches(5.55), Inches(11.9), Inches(1.2),
              "Neither limitation has been tested outside physics simulation. Discrete-event / queueing "
              "systems are a fundamentally different domain — stochastic discrete arrivals and departures, "
              "resource contention, priority scheduling — not continuous PDE fields. This is an open gap.",
              size=16, color=TEXT_DARK)

    # ---------------- Slide 4: Our Bridging Approach ----------------
    s = blank_slide(prs)
    slide_header(s, "Approach", "Our Bridging Approach", 4)
    add_text(s, Inches(0.7), Inches(1.65), Inches(11.9), Inches(0.7),
              "Apply Conformal Prediction — specifically Mondrian CP — to an ER queueing surrogate model, "
              "and directly test whether Gopakumar et al.'s two stated limitations hold or break here.",
              size=17, bold=True, color=DARK)

    col_w = Inches(5.7)
    add_rect(s, Inches(0.7), Inches(2.6), col_w, Inches(3.9), WHITE, line=True)
    add_rect(s, Inches(0.7), Inches(2.6), col_w, Inches(0.55), ACCENT_DARK)
    add_text(s, Inches(0.95), Inches(2.72), col_w - Inches(0.5), Inches(0.4),
              "Addressing: Marginal coverage", size=15, bold=True, color=WHITE)
    add_bullets(s, Inches(0.95), Inches(3.35), col_w - Inches(0.5), Inches(3.0), [
        "Mondrian CP partitions calibration by category (staffing level, shift, arrival-rate regime) "
        "instead of pooling everything into one marginal guarantee.",
        "Target: near-conditional coverage — each category individually close to the nominal rate, "
        "not just the average across all of them.",
    ], size=14, space_after=12)

    add_rect(s, Inches(6.7), Inches(2.6), col_w, Inches(3.9), WHITE, line=True)
    add_rect(s, Inches(6.7), Inches(2.6), col_w, Inches(0.55), ACCENT_DARK)
    add_text(s, Inches(6.95), Inches(2.72), col_w - Inches(0.5), Inches(0.4),
              "Addressing: Exchangeability", size=15, bold=True, color=WHITE)
    add_bullets(s, Inches(6.95), Inches(3.35), col_w - Inches(0.5), Inches(3.0), [
        "Phase 2 stress-test: simulate an out-of-distribution “surge day” scenario outside the "
        "training regime.",
        "Directly probes where standard CP and Mondrian CP hold or break when the exchangeability "
        "assumption is violated.",
    ], size=14, space_after=12)

    # ---------------- Slide 5: Methodology Overview ----------------
    s = blank_slide(prs)
    slide_header(s, "Pipeline", "Methodology Overview", 5)

    y = Inches(3.0)
    h = Inches(1.6)
    bw = Inches(3.15)
    aw = Inches(0.55)
    x = Inches(0.55)

    pipeline_box(s, x, y, bw, h, "Discrete-Event\nSimulation", "SimPy, calibrated on real ED data", ACCENT_DARK)
    x2 = x + bw
    pipeline_arrow(s, x2, y + h / 2 - Inches(0.2), aw, Inches(0.4))
    x3 = x2 + aw
    pipeline_box(s, x3, y, bw, h, "Surrogate Model", "Gradient boosting, trained on\nDES scenario sweep", ACCENT)
    x4 = x3 + bw
    pipeline_arrow(s, x4, y + h / 2 - Inches(0.2), aw, Inches(0.4))
    x5 = x4 + aw
    pipeline_box(s, x5, y, bw, h, "Uncertainty\nQuantification", "GP baseline vs. Standard CP\nvs. Mondrian CP", RGBColor(0x0F, 0x76, 0x6E))

    add_text(s, Inches(0.55), Inches(4.95), Inches(12.2), Inches(0.4),
              "Each stage validated against the previous before moving on:", size=14, italic=True, color=TEXT_MUTED)
    add_bullets(s, Inches(0.7), Inches(5.4), Inches(11.7), Inches(1.6), [
        "DES output validated against real aggregated ED statistics (91.0% match to real daily volume).",
        "Surrogate accuracy validated against DES outputs held out from training (MAE / RMSE / R²).",
        "UQ methods validated against a shared, fixed test split and shared target coverage (α = 0.1) "
        "so GP, standard CP, and Mondrian CP are directly comparable.",
    ], size=15, space_after=8)

    # ---------------- Slide 6: Real-world data used ----------------
    s = blank_slide(prs)
    slide_header(s, "Data", "Real-World Data Used", 6)
    add_text(s, Inches(0.7), Inches(1.6), Inches(11.9), Inches(0.5),
              "Hospital Triage and Patient History Data (Kaggle, maalona)", size=18, bold=True, color=DARK)
    add_text(s, Inches(0.7), Inches(2.05), Inches(11.9), Inches(0.5),
              "Yale New Haven Health System, retrospective study, March 2014 – July 2017",
              size=14, color=TEXT_MUTED, italic=True)

    stats = [
        ("560,486", "total ED visits"),
        ("972", "variables per visit"),
        ("3", "EDs (1 academic + 2 community)"),
        ("258.2 / day", "calibrated rate, Dept. A"),
    ]
    sx = Inches(0.7)
    sw = Inches(2.85)
    for i, (num, lbl) in enumerate(stats):
        cx = sx + i * (sw + Inches(0.1))
        add_rect(s, cx, Inches(2.75), sw, Inches(1.15), LIGHT_BG, line=True)
        add_text(s, cx, Inches(2.85), sw, Inches(0.55), num, size=22, bold=True, color=ACCENT_DARK, align=PP_ALIGN.CENTER)
        add_text(s, cx, Inches(3.4), sw, Inches(0.4), lbl, size=11, color=TEXT_MUTED, align=PP_ALIGN.CENTER)

    add_text(s, Inches(0.7), Inches(4.25), Inches(11.9), Inches(0.4),
              "What we extracted — real vs. literature-calibrated (disclosed explicitly, not blended silently):",
              size=15, bold=True, color=TEXT_DARK)

    add_rect(s, Inches(0.7), Inches(4.75), Inches(5.8), Inches(2.0), WHITE, line=True)
    add_text(s, Inches(0.9), Inches(4.85), Inches(5.4), Inches(0.35), "From real data", size=13, bold=True, color=GOOD)
    add_bullets(s, Inches(0.9), Inches(5.2), Inches(5.4), Inches(1.5), [
        "Arrival pattern by 4-hour bin, day of week, month",
        "ESI acuity mix (triage severity distribution)",
        "Department-level daily visit rate",
    ], size=13, space_after=6)

    add_rect(s, Inches(6.65), Inches(4.75), Inches(5.95), Inches(2.0), WHITE, line=True)
    add_text(s, Inches(6.85), Inches(4.85), Inches(5.4), Inches(0.35), "From literature (not in dataset)", size=13, bold=True, color=WARN)
    add_bullets(s, Inches(6.85), Inches(5.2), Inches(5.4), Inches(1.5), [
        "Service/treatment time per ESI level (log-normal)",
        "Dataset has no discharge timestamp — confirmed by",
        ("exhaustive search across all 972 columns", 1),
    ], size=13, space_after=6)

    # ---------------- Slide 7: Work Completed ----------------
    s = blank_slide(prs)
    slide_header(s, "Progress — 50%", "Work Completed", 7)

    add_text(s, Inches(0.7), Inches(1.55), Inches(11.9), Inches(0.4),
              "DES validated against real data (200 simulated days)", size=15, bold=True, color=DARK)
    add_table(s, Inches(0.7), Inches(1.95), Inches(5.6), Inches(1.1),
              ["Metric", "Value"],
              [
                  ["Real visits/day (Dept. A)", "258.2"],
                  ["Simulated mean visits/day", "235.1"],
                  ["Match", "91.0%"],
              ],
              col_widths=[Inches(3.4), Inches(2.2)], font_size=13)

    add_text(s, Inches(6.9), Inches(1.55), Inches(5.7), Inches(0.4),
              "Surrogate model accuracy (gradient boosting)", size=15, bold=True, color=DARK)
    add_table(s, Inches(6.9), Inches(1.95), Inches(5.7), Inches(1.55),
              ["Target", "MAE", "RMSE", "R²"],
              [
                  ["n_patients", "9.90", "12.59", "0.929"],
                  ["mean_wait_min", "8.86", "13.23", "0.787"],
                  ["mean_total_min", "9.88", "13.61", "0.762"],
                  ["p95_wait_min", "66.94", "102.47", "0.647"],
              ],
              col_widths=[Inches(2.0), Inches(1.2), Inches(1.2), Inches(1.3)], font_size=12)

    add_text(s, Inches(0.7), Inches(3.65), Inches(11.9), Inches(0.4),
              "GP baseline UQ (α = 0.1, i.e. 90% target coverage)", size=15, bold=True, color=DARK)
    add_table(s, Inches(0.7), Inches(4.05), Inches(11.9), Inches(1.7),
              ["Target", "Target coverage", "Empirical coverage", "Mean interval width"],
              [
                  ["n_patients", "90%", "88.5%", "39.3"],
                  ["mean_wait_minutes", "90%", "87.7%", "43.6"],
                  ["mean_total_minutes", "90%", "88.8%", "43.4"],
                  ["p95_wait_minutes", "90%", "90.1%", "350.3"],
              ],
              col_widths=[Inches(3.5), Inches(2.8), Inches(2.8), Inches(2.8)], font_size=13)

    add_text(s, Inches(0.7), Inches(6.0), Inches(11.9), Inches(1.0),
              "GP undercovers on 3/4 targets (87.7–88.8% vs. 90% target) — expected, since GP relies on "
              "distributional assumptions rather than a finite-sample guarantee. This is precisely the gap "
              "we test whether CP / Mondrian CP close in Phase 2.",
              size=13, italic=True, color=TEXT_MUTED)

    # ---------------- Slide 8: Literature review snapshot (placeholder) ----------------
    s = blank_slide(prs)
    slide_header(s, "Background", "Literature Review Snapshot", 8)
    add_rect(s, Inches(0.7), Inches(1.7), Inches(11.9), Inches(0.6), RGBColor(0xFE, 0xF3, 0xC7), line=True)
    add_text(s, Inches(0.95), Inches(1.83), Inches(11.4), Inches(0.4),
              "PLACEHOLDER — literature review in progress, update before presenting.", size=13, bold=True, color=WARN)

    cats = ["Conformal Prediction\nfoundations", "Surrogate modeling", "Queueing theory /\nED simulation"]
    cw = Inches(3.75)
    for i, cat in enumerate(cats):
        cx = Inches(0.7) + i * (cw + Inches(0.15))
        add_rect(s, cx, Inches(2.7), cw, Inches(2.4), LIGHT_BG, line=True)
        add_text(s, cx, Inches(2.9), cw, Inches(0.8), cat, size=16, bold=True, color=DARK, align=PP_ALIGN.CENTER)
        add_text(s, cx, Inches(3.9), cw, Inches(0.6), "[ __ ] papers reviewed", size=15, color=TEXT_MUTED, align=PP_ALIGN.CENTER)

    add_text(s, Inches(0.7), Inches(5.5), Inches(11.9), Inches(0.5),
              "Target: 30 papers total, 3 core papers reviewed in depth (incl. Gopakumar et al. 2026).",
              size=15, color=TEXT_DARK)

    # ---------------- Slide 9: Future scope ----------------
    s = blank_slide(prs)
    slide_header(s, "Roadmap", "Future Scope — End-Sem Preview", 9)
    items = [
        ("Week 9–10", "Implement standard conformal prediction on surrogate residuals; test multiple nonconformity measures."),
        ("Week 11–12", "Implement Mondrian CP — partition by staffing level / shift / arrival-rate category; measure per-category coverage."),
        ("Week 13", "Stress-test exchangeability with an out-of-distribution “surge day” scenario; check where standard CP and Mondrian CP hold or break."),
        ("Week 14–15", "Full comparison — GP vs. standard CP vs. Mondrian CP: coverage, interval width, computation time."),
        ("Week 16", "End-sem PPT: quantified answer on whether Mondrian CP closes the marginal-vs-conditional coverage gap in this domain; final report."),
    ]
    y = Inches(1.7)
    for wk, desc in items:
        add_rect(s, Inches(0.7), y, Inches(1.7), Inches(0.85), ACCENT_DARK)
        add_text(s, Inches(0.7), y, Inches(1.7), Inches(0.85), wk, size=13, bold=True, color=WHITE,
                  align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, Inches(2.55), y, Inches(10.0), Inches(0.85), desc, size=13.5, color=TEXT_DARK,
                  anchor=MSO_ANCHOR.MIDDLE)
        y += Inches(1.0)

    prs.save(OUT_PATH)
    print(f"Saved {OUT_PATH}")


if __name__ == "__main__":
    build()
