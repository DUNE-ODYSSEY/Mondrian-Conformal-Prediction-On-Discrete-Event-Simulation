"""
Zeroth review PPT - exact structure per professor's template: title (names,
roll numbers, group number), table of contents, abstract, literature review
(table format), research gap, problem statement, methodology (diagrams,
multiple slides), results, conclusion, references, thank you.

Literature review and references list only real, verified items (the base
paper, the dataset) - no fabricated citations, since the literature review
itself is still in progress. Results slide shows real preliminary numbers
already produced (DES validation, surrogate accuracy), not the full
publication-grade findings from later in the project - proportionate to a
zeroth review.

Re-run: .venv\\Scripts\\python.exe slides\\build_zeroth_review_ppt.py

Same polished custom-drawn design as build_full_results_ppt.py (user's
stated preference for this project's decks).
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

OUT_PATH = "slides/zeroth_review.pptx"

TEAM = [
    ("G Venugopalan", "CB.AI.U4AID25115"),
    ("Vipin Sudhakar", "CB.AI.U4AID25166"),
    ("Rithvik Arulprakash", "CB.AI.U4AID25148"),
    ("Harshith Kv", "CB.AI.U4AID25119"),
]
COURSE_CODE = "23AID201"
GROUP_NO = "B9"
PROJECT_TITLE = "Mondrian Conformal Prediction for Uncertainty\nQuantification in ER Discrete-Event Simulation Surrogates"

DARK = RGBColor(0x1B, 0x22, 0x38)
ACCENT = RGBColor(0x3B, 0x82, 0xF6)
ACCENT_DARK = RGBColor(0x1E, 0x40, 0xAF)
TEXT_DARK = RGBColor(0x1F, 0x29, 0x37)
TEXT_MUTED = RGBColor(0x6B, 0x72, 0x80)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_BG = RGBColor(0xF8, 0xFA, 0xFC)
BORDER = RGBColor(0xE2, 0xE8, 0xF0)
GOOD = RGBColor(0x16, 0xA3, 0x4A)
WARN = RGBColor(0xD9, 0x77, 0x06)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
TOTAL_SLIDES = 12


def new_deck():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


def blank_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def add_rect(slide, x, y, w, h, color, line=False):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    if line:
        shape.line.color.rgb = BORDER
        shape.line.width = Pt(0.75)
    else:
        shape.line.fill.background()
    shape.shadow.inherit = False
    return shape


def add_text(slide, x, y, w, h, text, size=18, bold=False, color=TEXT_DARK,
             align=PP_ALIGN.LEFT, font="Calibri", anchor=MSO_ANCHOR.TOP, italic=False):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    lines = text.split("\n")
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = line
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.italic = italic
        run.font.color.rgb = color
        run.font.name = font
    return box


def add_bullets(slide, x, y, w, h, items, size=16, color=TEXT_DARK, space_after=10, font="Calibri"):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(space_after)
        text, level = item if isinstance(item, tuple) else (item, 0)
        p.level = level
        bullet = "•  " if level == 0 else "‒  "
        run = p.add_run()
        run.text = bullet + text
        run.font.size = Pt(size - level * 2)
        run.font.color.rgb = color
        run.font.name = font
    return box


def slide_header(slide, kicker, title, num, total=TOTAL_SLIDES, title_size=28):
    add_rect(slide, 0, 0, SLIDE_W, Inches(0.09), ACCENT)
    add_text(slide, Inches(0.6), Inches(0.3), Inches(11.5), Inches(0.35),
              kicker.upper(), size=13, bold=True, color=ACCENT, font="Calibri")
    add_text(slide, Inches(0.6), Inches(0.62), Inches(11.9), Inches(0.7),
              title, size=title_size, bold=True, color=DARK, font="Calibri")
    add_rect(slide, Inches(0.6), Inches(1.28), Inches(1.1), Pt(3), ACCENT)
    add_text(slide, Inches(12.2), Inches(7.05), Inches(0.9), Inches(0.35),
              f"{num} / {total}", size=11, color=TEXT_MUTED, align=PP_ALIGN.RIGHT)


def add_table(slide, x, y, w, h, headers, rows, col_widths=None, font_size=13, header_color=ACCENT_DARK):
    n_rows = len(rows) + 1
    n_cols = len(headers)
    gtable = slide.shapes.add_table(n_rows, n_cols, x, y, w, h).table
    if col_widths:
        for i, cw in enumerate(col_widths):
            gtable.columns[i].width = cw
    for j, hdr in enumerate(headers):
        cell = gtable.cell(0, j)
        cell.text = hdr
        cell.fill.solid()
        cell.fill.fore_color.rgb = header_color
        for p in cell.text_frame.paragraphs:
            p.alignment = PP_ALIGN.CENTER if j > 0 else PP_ALIGN.LEFT
            for run in p.runs:
                run.font.size = Pt(font_size)
                run.font.bold = True
                run.font.color.rgb = WHITE
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        cell.margin_top = Pt(4)
        cell.margin_bottom = Pt(4)
    for i, row in enumerate(rows):
        for j, val in enumerate(row):
            cell = gtable.cell(i + 1, j)
            cell.text = str(val)
            cell.fill.solid()
            cell.fill.fore_color.rgb = LIGHT_BG if i % 2 == 0 else WHITE
            # cell.text splits on "\n" into separate paragraphs, not one run with
            # a line break — style every paragraph/run, not just paragraphs[0],
            # or unstyled lines fall back to a much larger default font size.
            for p in cell.text_frame.paragraphs:
                p.alignment = PP_ALIGN.CENTER if j > 0 else PP_ALIGN.LEFT
                for run in p.runs:
                    run.font.size = Pt(font_size)
                    run.font.color.rgb = TEXT_DARK
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            cell.margin_top = Pt(4)
            cell.margin_bottom = Pt(4)
    return gtable


def build():
    prs = new_deck()
    n = [0]

    def nn():
        n[0] += 1
        return n[0]

    # ---------------- Slide 1: Title ----------------
    s = blank_slide(prs)
    add_rect(s, 0, 0, SLIDE_W, SLIDE_H, DARK)
    add_rect(s, 0, Inches(6.55), SLIDE_W, Inches(0.12), ACCENT)
    add_text(s, Inches(1), Inches(1.05), Inches(11.3), Inches(2.0),
              PROJECT_TITLE, size=30, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, Inches(1), Inches(3.05), Inches(11.3), Inches(0.5),
              "Zeroth Review", size=17, color=RGBColor(0x93, 0xC5, 0xFD), align=PP_ALIGN.CENTER, italic=True)
    add_text(s, Inches(1), Inches(3.55), Inches(11.3), Inches(0.4),
              f"Group: {GROUP_NO}", size=15, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    y = Inches(4.35)
    for nm, roll in TEAM:
        add_text(s, Inches(1), y, Inches(11.3), Inches(0.4), f"{nm}   —   {roll}",
                  size=14, color=WHITE, align=PP_ALIGN.CENTER)
        y += Inches(0.42)

    add_text(s, Inches(1), Inches(6.15), Inches(11.3), Inches(0.4),
              f"Course: {COURSE_CODE}", size=13, color=RGBColor(0x93, 0xC5, 0xFD), align=PP_ALIGN.CENTER)
    add_text(s, Inches(12.2), Inches(7.05), Inches(0.9), Inches(0.35),
              f"{nn()} / {TOTAL_SLIDES}", size=11, color=RGBColor(0x93, 0xC5, 0xFD), align=PP_ALIGN.RIGHT)

    # ---------------- Slide 2: Table of Contents ----------------
    s = blank_slide(prs)
    slide_header(s, "Outline", "Table of Contents", nn())
    toc_items = [
        ("01", "Abstract"),
        ("02", "Literature Review"),
        ("03", "Research Gap"),
        ("04", "Problem Statement"),
        ("05", "Methodology"),
        ("06", "Results"),
        ("07", "Conclusion"),
        ("08", "References"),
    ]
    y = Inches(1.75)
    row_h = Inches(0.62)
    for num_label, title in toc_items:
        add_text(s, Inches(0.9), y, Inches(0.9), row_h, num_label, size=18, bold=True, color=ACCENT_DARK)
        add_text(s, Inches(2.0), y, Inches(9.5), row_h, title, size=17, color=TEXT_DARK, anchor=MSO_ANCHOR.MIDDLE)
        add_rect(s, Inches(0.9), y + row_h - Pt(2), Inches(11.4), Pt(1), BORDER)
        y += row_h

    # ---------------- Slide 3: Abstract ----------------
    s = blank_slide(prs)
    slide_header(s, "Summary", "Abstract", nn())
    add_rect(s, Inches(0.7), Inches(1.7), Inches(11.9), Inches(4.8), LIGHT_BG, line=True)
    add_text(s, Inches(1.0), Inches(2.0), Inches(11.3), Inches(4.2),
              "Surrogate models are increasingly used to approximate expensive simulations and guide "
              "decisions in stochastic service systems such as emergency room staffing. However, reliable "
              "uncertainty quantification (UQ) for these surrogates remains an open problem: Gaussian "
              "Processes are costly and rely on distributional assumptions, while Conformal Prediction (CP) "
              "offers a distribution-free, finite-sample coverage guarantee but has only been validated for "
              "surrogate models in physics domains (PDEs, magnetohydrodynamics, weather, fusion).\n\n"
              "This project builds a discrete-event simulation (DES) of an emergency room, calibrated on "
              "real hospital data, trains a surrogate model on it, and applies standard and Mondrian "
              "Conformal Prediction to quantify uncertainty in the surrogate's predictions. The goal is to "
              "test whether two limitations explicitly noted in prior work — marginal-only coverage and an "
              "untested exchangeability assumption — hold up in a discrete-event / queueing domain that has "
              "not been studied before in this context.",
              size=16, color=TEXT_DARK)

    # ---------------- Slide 4: Literature Review (table format) ----------------
    s = blank_slide(prs)
    slide_header(s, "Background", "Literature Review", nn())
    add_table(s, Inches(0.6), Inches(1.65), Inches(12.1), Inches(3.6),
              ["S.No", "Focus Area", "Key Reference / Topic", "Relevance to Our Work"],
              [
                  ["1", "Base Paper", "Gopakumar et al. (2026) — CP for\nsurrogate UQ in physics domains",
                   "States the two limitations (marginal\ncoverage, exchangeability) this project tests"],
                  ["2", "Conformal Prediction\nFoundations", "Split conformal prediction;\nMondrian (category-conditional) CP",
                   "Core UQ methodology used\nin this project"],
                  ["3", "Surrogate Modeling", "Gradient boosting / GP / neural-\nnetwork surrogates for simulations",
                   "Basis for our DES\nsurrogate model"],
                  ["4", "Queueing Theory /\nDiscrete-Event Simulation", "Discrete-event simulation of\nED / service systems",
                   "Basis for our ER\nsimulation design"],
              ],
              col_widths=[Inches(0.9), Inches(3.0), Inches(4.1), Inches(4.1)], font_size=12)
    add_text(s, Inches(0.6), Inches(5.55), Inches(12.1), Inches(0.5),
              "Full review in progress — target 30 papers total, 3 core papers reviewed in depth.",
              size=13, italic=True, color=TEXT_MUTED)

    # ---------------- Slide 5: Research Gap ----------------
    s = blank_slide(prs)
    slide_header(s, "Gaps Identified", "Research Gap", nn())
    add_text(s, Inches(0.7), Inches(1.65), Inches(11.9), Inches(0.5),
              "Gopakumar et al. (2026) validate Conformal Prediction for surrogate-model UQ —", size=17, bold=True)
    add_text(s, Inches(0.7), Inches(2.08), Inches(11.9), Inches(0.5),
              "but only in physics domains: PDEs, magnetohydrodynamics (MHD), weather, fusion.", size=17, color=ACCENT_DARK, bold=True)

    gaps = [
        ("Gap 1 — Marginal coverage only", "No guarantee of coverage within specific subgroups or "
         "conditions, only on average across everything."),
        ("Gap 2 — Untested exchangeability assumption", "Coverage relies on calibration and test data being "
         "exchangeable; this hasn't been tested under distribution shift."),
        ("Gap 3 — No discrete-event / queueing validation", "Neither limitation has been tested outside "
         "physics simulation — discrete-event / queueing systems are a fundamentally different domain."),
    ]
    y = Inches(2.75)
    for title, desc in gaps:
        add_rect(s, Inches(0.7), y, Inches(11.9), Inches(1.25), LIGHT_BG, line=True)
        add_text(s, Inches(1.0), y + Inches(0.12), Inches(11.3), Inches(0.4), title, size=15, bold=True, color=ACCENT_DARK)
        add_text(s, Inches(1.0), y + Inches(0.55), Inches(11.3), Inches(0.6), desc, size=13.5, color=TEXT_DARK)
        y += Inches(1.4)

    # ---------------- Slide 6: Problem Statement ----------------
    s = blank_slide(prs)
    slide_header(s, "Motivation", "Problem Statement", nn())
    add_bullets(s, Inches(0.7), Inches(1.75), Inches(11.9), Inches(5), [
        "Surrogate models are fast, learned approximations of expensive simulations, increasingly used "
        "to guide decisions in stochastic service systems - ER staffing, queueing networks, etc.",
        "A point prediction alone isn't enough for a decision like \"will adding 2 more doctors bring "
        "wait time down enough?\" - that needs a calibrated interval, not just a number.",
        "Gaussian Processes are the established way to get one, but are expensive to train and rely on "
        "distributional assumptions rather than a finite-sample guarantee.",
        "Conformal Prediction (CP) is distribution-free with a finite-sample coverage guarantee and is "
        "cheap to calibrate - but it's only been validated for surrogate models in a narrow set of "
        "domains so far.",
    ], size=18, space_after=16)

    # ---------------- Slide 7: Methodology I — Pipeline ----------------
    s = blank_slide(prs)
    slide_header(s, "Approach (1/2)", "Methodology — Pipeline Overview", nn(), title_size=25)
    y = Inches(2.1)
    h = Inches(1.3)
    bw = Inches(3.6)
    gap = Inches(0.45)
    x_positions = [Inches(0.55), Inches(0.55) + bw + gap, Inches(0.55) + 2 * (bw + gap)]
    labels = [
        ("Discrete-Event\nSimulation", "SimPy, calibrated on\nreal ED data"),
        ("Surrogate Model", "Learned approximation of\nthe simulation"),
        ("Uncertainty\nQuantification", "GP baseline vs. standard CP\nvs. Mondrian CP"),
    ]
    for x, (label, sub) in zip(x_positions, labels):
        box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, bw, h)
        box.adjustments[0] = 0.08
        box.fill.solid()
        box.fill.fore_color.rgb = ACCENT_DARK
        box.line.fill.background()
        box.shadow.inherit = False
        tf = box.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = label
        r.font.size = Pt(16)
        r.font.bold = True
        r.font.color.rgb = WHITE
        p2 = tf.add_paragraph()
        p2.alignment = PP_ALIGN.CENTER
        r2 = p2.add_run()
        r2.text = sub
        r2.font.size = Pt(11.5)
        r2.font.color.rgb = WHITE
    for x in x_positions[:-1]:
        arrow = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, x + bw, y + h / 2 - Inches(0.2), gap, Inches(0.4))
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = TEXT_MUTED
        arrow.line.fill.background()
        arrow.shadow.inherit = False

    add_text(s, Inches(0.55), Inches(4.0), Inches(12.2), Inches(0.4),
              "Each stage validated against the previous before moving on:", size=14, italic=True, color=TEXT_MUTED)
    add_bullets(s, Inches(0.7), Inches(4.45), Inches(11.7), Inches(2.4), [
        "DES output validated against real aggregated ED statistics before being trusted as a data source.",
        "Surrogate accuracy validated against DES outputs held out from training.",
        "Every UQ method evaluated on the same test data and the same target coverage level, so results "
        "are directly comparable across methods.",
    ], size=15, space_after=10)

    # ---------------- Slide 8: Methodology II — UQ approaches compared ----------------
    s = blank_slide(prs)
    slide_header(s, "Approach (2/2)", "Methodology — UQ Methods Compared", nn(), title_size=24)
    cards = [
        ("GP Baseline", "Gaussian Process regression gives a predictive mean and standard deviation; "
         "interval = mean ± z·std. Established, but costly to train and relies on distributional "
         "assumptions.", ACCENT_DARK),
        ("Standard CP", "Calibrates a single interval width from the surrogate's residuals on held-out "
         "calibration data. Distribution-free, finite-sample coverage guarantee — but only on average "
         "across all conditions.", ACCENT_DARK),
        ("Mondrian CP", "Calibrates a separate interval width per category (e.g. staffing level x "
         "arrival-rate regime) instead of one pooled value — targets conditional coverage, addressing "
         "Gap 1 directly.", RGBColor(0x0F, 0x76, 0x6E)),
    ]
    cw = Inches(3.85)
    for i, (title, desc, color) in enumerate(cards):
        cx = Inches(0.6) + i * (cw + Inches(0.15))
        add_rect(s, cx, Inches(1.75), cw, Inches(4.6), WHITE, line=True)
        add_rect(s, cx, Inches(1.75), cw, Inches(0.7), color)
        add_text(s, cx + Inches(0.15), Inches(1.85), cw - Inches(0.3), Inches(0.5),
                  title, size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_text(s, cx + Inches(0.25), Inches(2.65), cw - Inches(0.5), Inches(3.5),
                  desc, size=13.5, color=TEXT_DARK)

    # ---------------- Slide 9: Results (preliminary) ----------------
    s = blank_slide(prs)
    slide_header(s, "Preliminary", "Results", nn())
    add_table(s, Inches(0.7), Inches(1.7), Inches(11.9), Inches(1.5),
              ["DES Validation (200 simulated days)", "Value"],
              [
                  ["Real visits/day (Dept. A)", "258.2"],
                  ["Simulated mean visits/day", "235.1"],
                  ["Match to real data", "91.0%"],
              ], col_widths=[Inches(8.5), Inches(3.4)], font_size=14)

    add_table(s, Inches(0.7), Inches(3.55), Inches(11.9), Inches(2.0),
              ["Surrogate Target", "MAE", "RMSE", "R²"],
              [
                  ["n_patients", "9.90", "12.59", "0.929"],
                  ["mean_wait_minutes", "8.86", "13.23", "0.787"],
                  ["mean_total_minutes", "9.88", "13.61", "0.762"],
                  ["p95_wait_minutes", "66.94", "102.47", "0.647"],
              ], col_widths=[Inches(3.9), Inches(2.7), Inches(2.7), Inches(2.6)], font_size=13)

    add_text(s, Inches(0.7), Inches(5.85), Inches(11.9), Inches(0.9),
              "The simulation is validated against real hospital data and the surrogate model fits well "
              "on 3 of 4 targets. Full evaluation of GP baseline, standard CP, and Mondrian CP coverage "
              "is in progress.", size=14, italic=True, color=TEXT_MUTED)

    # ---------------- Slide 10: Conclusion ----------------
    s = blank_slide(prs)
    slide_header(s, "Summary", "Conclusion", nn())
    add_bullets(s, Inches(0.7), Inches(1.75), Inches(11.9), Inches(4), [
        "This project tests two limitations explicitly flagged by Gopakumar et al. (2026) — marginal "
        "coverage and the exchangeability assumption — in a discrete-event / queueing domain that has "
        "not been studied before in this context.",
        "The core contribution is applying Mondrian Conformal Prediction to close the marginal-vs-"
        "conditional coverage gap, and stress-testing exchangeability with an out-of-distribution "
        "scenario, using a real hospital-calibrated ER simulation as the testbed.",
        "Simulation and surrogate model are built and validated; uncertainty quantification methods "
        "(GP baseline, standard CP, Mondrian CP) are the focus of the next phase.",
    ], size=17, space_after=16)

    # ---------------- Slide 11: References ----------------
    s = blank_slide(prs)
    slide_header(s, "Sources", "References", nn())
    add_bullets(s, Inches(0.7), Inches(1.8), Inches(11.9), Inches(3), [
        "Gopakumar, et al. (2026). Conformal Prediction for surrogate-model uncertainty quantification "
        "in physics simulation domains (PDEs, MHD, weather, fusion). [Base paper]",
        "Hospital Triage and Patient History Data. Kaggle (maalona) — Yale New Haven Health System, "
        "retrospective study, March 2014 – July 2017. [Real-world calibration data]",
    ], size=16, space_after=16)
    add_rect(s, Inches(0.7), Inches(4.0), Inches(11.9), Inches(0.8), LIGHT_BG, line=True)
    add_text(s, Inches(1.0), Inches(4.2), Inches(11.3), Inches(0.4),
              "Full reference list (30 papers target) pending completion of the literature review.",
              size=13, italic=True, color=WARN)

    # ---------------- Slide 12: Thank You ----------------
    s = blank_slide(prs)
    add_rect(s, 0, 0, SLIDE_W, SLIDE_H, DARK)
    add_rect(s, 0, Inches(6.55), SLIDE_W, Inches(0.12), ACCENT)
    add_text(s, Inches(1), Inches(3.0), Inches(11.3), Inches(1.0),
              "Thank You", size=44, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, Inches(1), Inches(4.0), Inches(11.3), Inches(0.5),
              "Questions?", size=18, color=RGBColor(0x93, 0xC5, 0xFD), align=PP_ALIGN.CENTER, italic=True)
    add_text(s, Inches(12.2), Inches(7.05), Inches(0.9), Inches(0.35),
              f"{nn()} / {TOTAL_SLIDES}", size=11, color=RGBColor(0x93, 0xC5, 0xFD), align=PP_ALIGN.RIGHT)

    prs.save(OUT_PATH)
    print(f"Saved {OUT_PATH}")


if __name__ == "__main__":
    build()
